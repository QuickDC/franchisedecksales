#!/usr/bin/env python3
"""
Generate a premium, keynote-style PowerPoint presentation
for a dry cleaning business seminar in Jodhpur.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# === COLORS ===
DARK_NAVY = RGBColor(10, 25, 47)
WHITE = RGBColor(255, 255, 255)
MUTED_ORANGE = RGBColor(235, 120, 60)
GOLD_ACCENT = RGBColor(200, 160, 80)
LIGHT_GRAY = RGBColor(240, 240, 245)
MEDIUM_GRAY = RGBColor(100, 100, 110)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs

def set_slide_background(slide, color=DARK_NAVY):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_title_slide(prs, title, subtitle, footer_text):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)

    # Decorative orange accent bar at top
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = MUTED_ORANGE
    shape.line.fill.background()

    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), prs.slide_width - Inches(1), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), prs.slide_width - Inches(2), Inches(1.2))
    tf = sub_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 210)
    p.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), prs.slide_width - Inches(1), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = footer_text
    p.font.size = Pt(18)
    p.font.color.rgb = MUTED_ORANGE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, points, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)

    # Orange accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MUTED_ORANGE
    shape.line.fill.background()

    # Section label if provided
    if section:
        section_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.4))
        tf = section_box.text_frame
        p = tf.paragraphs[0]
        p.text = section.upper()
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED_ORANGE
        p.font.bold = True
        p.letter_spacing = Pt(2)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), prs.slide_width - Inches(1.6), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Points
    points_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), prs.slide_width - Inches(1.6), Inches(5))
    tf = points_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(22)
        p.font.color.rgb = WHITE
        p.space_after = Pt(16)

    return slide

def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)

    # Orange accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MUTED_ORANGE
    shape.line.fill.background()

    # Section label if provided
    if section:
        section_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.4))
        tf = section_box.text_frame
        p = tf.paragraphs[0]
        p.text = section.upper()
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED_ORANGE
        p.font.bold = True
        p.letter_spacing = Pt(2)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.8), prs.slide_width - Inches(1.6), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(6), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = MUTED_ORANGE

    # Left points
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.4), Inches(6), Inches(4.5))
    tf = left_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(left_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(7), Inches(1.8), Inches(6), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = GOLD_ACCENT

    # Right points
    right_box = slide.shapes.add_textbox(Inches(7), Inches(2.4), Inches(6), Inches(4.5))
    tf = right_box.text_frame
    tf.word_wrap = True

    for i, point in enumerate(right_points):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(18)
        p.font.color.rgb = WHITE
        p.space_after = Pt(12)

    return slide

def add_framework_slide(prs, title, items, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)

    # Orange accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MUTED_ORANGE
    shape.line.fill.background()

    # Section label if provided
    if section:
        section_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.4))
        tf = section_box.text_frame
        p = tf.paragraphs[0]
        p.text = section.upper()
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED_ORANGE
        p.font.bold = True
        p.letter_spacing = Pt(2)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), prs.slide_width - Inches(1.6), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Framework cards
    card_width = Inches(2.3)
    card_height = Inches(3.5)
    start_x = Inches(0.8)
    gap = Inches(0.3)
    y_pos = Inches(1.8)

    colors = [MUTED_ORANGE, GOLD_ACCENT, RGBColor(80, 130, 180), RGBColor(100, 160, 100), RGBColor(150, 100, 160), RGBColor(180, 100, 80)]

    for i, item in enumerate(items):
        x = start_x + i * (card_width + gap)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = colors[i % len(colors)]
        card.line.fill.background()

        # Number
        num_box = slide.shapes.add_textbox(x + Inches(0.2), y_pos + Inches(0.3), card_width - Inches(0.4), Inches(0.6))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Title
        title_p = slide.shapes.add_textbox(x + Inches(0.2), y_pos + Inches(1.1), card_width - Inches(0.4), Inches(0.6))
        tf = title_p.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item['title']
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Description
        desc = slide.shapes.add_textbox(x + Inches(0.2), y_pos + Inches(1.8), card_width - Inches(0.4), Inches(1.5))
        tf = desc.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item['desc']
        p.font.size = Pt(13)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_hexagon_slide(prs, title, items, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)

    # Orange accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MUTED_ORANGE
    shape.line.fill.background()

    # Section label if provided
    if section:
        section_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.4))
        tf = section_box.text_frame
        p = tf.paragraphs[0]
        p.text = section.upper()
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED_ORANGE
        p.font.bold = True
        p.letter_spacing = Pt(2)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), prs.slide_width - Inches(1.6), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Center circle
    center_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.5), Inches(3), Inches(2.33), Inches(2.33))
    center_circle.fill.solid()
    center_circle.fill.fore_color.rgb = MUTED_ORANGE
    center_circle.line.color.rgb = WHITE
    center_circle.line.width = Inches(0.05)

    center_text = slide.shapes.add_textbox(Inches(5.5), Inches(3.6), Inches(2.33), Inches(1))
    tf = center_text.text_frame
    p = tf.paragraphs[0]
    p.text = "GROWTH"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Outer hexagons
    positions = [
        (2.5, 1.8), (6.8, 1.8), (9.5, 3),
        (9.5, 5), (6.8, 6.2), (2.5, 5)
    ]

    colors = [MUTED_ORANGE, GOLD_ACCENT, RGBColor(80, 130, 180), RGBColor(100, 160, 100), RGBColor(150, 100, 160), RGBColor(180, 100, 80)]

    for i, item in enumerate(items):
        x = Inches(positions[i][0])
        y = Inches(positions[i][1])

        hex_shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, x, y, Inches(2), Inches(1.5))
        hex_shape.fill.solid()
        hex_shape.fill.fore_color.rgb = colors[i]
        hex_shape.line.fill.background()

        hex_text = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.4), Inches(1.8), Inches(0.8))
        tf = hex_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

    return slide

def add_cta_slide(prs, title, cta_items, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)

    # Orange accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MUTED_ORANGE
    shape.line.fill.background()

    # Section label if provided
    if section:
        section_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.4))
        tf = section_box.text_frame
        p = tf.paragraphs[0]
        p.text = section.upper()
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED_ORANGE
        p.font.bold = True
        p.letter_spacing = Pt(2)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), prs.slide_width - Inches(1.6), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # CTA cards
    card_width = Inches(3.5)
    card_height = Inches(2.8)
    start_x = Inches(1.5)
    gap = Inches(0.5)
    y_pos = Inches(2)

    for i, cta in enumerate(cta_items):
        x = start_x + i * (card_width + gap)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_pos, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY
        card.line.color.rgb = MUTED_ORANGE
        card.line.width = Inches(0.03)

        # Icon area (simulated with colored circle)
        icon_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.25), y_pos + Inches(0.3), Inches(1), Inches(1))
        icon_circle.fill.solid()
        icon_circle.fill.fore_color.rgb = MUTED_ORANGE
        icon_circle.line.fill.background()

        # CTA text
        cta_text = slide.shapes.add_textbox(x + Inches(0.3), y_pos + Inches(1.5), card_width - Inches(0.6), Inches(1))
        tf = cta_text.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cta
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = DARK_NAVY
        p.alignment = PP_ALIGN.CENTER

    # Footer note
    footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), prs.slide_width - Inches(1.6), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Visit our booth | Scan QR code | Talk to our team today"
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_outcome_slide(prs, title, outcomes, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)

    # Orange accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MUTED_ORANGE
    shape.line.fill.background()

    # Section label if provided
    if section:
        section_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.4))
        tf = section_box.text_frame
        p = tf.paragraphs[0]
        p.text = section.upper()
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED_ORANGE
        p.font.bold = True
        p.letter_spacing = Pt(2)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), prs.slide_width - Inches(1.6), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), prs.slide_width - Inches(1.6), Inches(0.5))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "QDC delivers measurable outcomes across every touchpoint"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(180, 180, 190)

    # Outcome items in 2 rows
    y_start = Inches(2.3)
    card_width = Inches(3.8)
    card_height = Inches(1.8)
    gap = Inches(0.4)

    for i, outcome in enumerate(outcomes):
        row = i // 3
        col = i % 3
        x = Inches(0.8 + col * (card_width + gap))
        y = y_start + row * (card_height + gap)

        # Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(30, 45, 67)
        card.line.fill.background()

        # Arrow icon
        arrow = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(0.5), Inches(0.4))
        tf = arrow.text_frame
        p = tf.paragraphs[0]
        p.text = "→"
        p.font.size = Pt(24)
        p.font.color.rgb = MUTED_ORANGE

        # Title
        title_p = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.6), card_width - Inches(0.4), Inches(0.5))
        tf = title_p.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = outcome['title']
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Description
        desc = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.1), card_width - Inches(0.4), Inches(0.6))
        tf = desc.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = outcome['desc']
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(180, 180, 190)

    return slide

def add_ecosystem_slide(prs, title, sections, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)

    # Orange accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MUTED_ORANGE
    shape.line.fill.background()

    # Section label if provided
    if section:
        section_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.4))
        tf = section_box.text_frame
        p = tf.paragraphs[0]
        p.text = section.upper()
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED_ORANGE
        p.font.bold = True
        p.letter_spacing = Pt(2)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), prs.slide_width - Inches(1.6), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(38)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), prs.slide_width - Inches(1.6), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "One connected platform powering every aspect of your business"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(180, 180, 190)

    # Section cards in 2 rows of 3
    y_start = Inches(1.9)
    card_width = Inches(3.8)
    card_height = Inches(2.2)
    gap = Inches(0.4)

    colors = [MUTED_ORANGE, GOLD_ACCENT, RGBColor(80, 130, 180), RGBColor(100, 160, 100), RGBColor(150, 100, 160), RGBColor(180, 100, 80)]

    for i, sec in enumerate(sections):
        row = i // 3
        col = i % 3
        x = Inches(0.8 + col * (card_width + gap))
        y = y_start + row * (card_height + gap)

        # Card background
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, card_width, card_height)
        card.fill.solid()
        card.fill.fore_color.rgb = colors[i]
        card.line.fill.background()

        # Section name
        name_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.3), card_width - Inches(0.6), Inches(0.5))
        tf = name_box.text_frame
        p = tf.paragraphs[0]
        p.text = sec['title']
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # Features
        features_box = slide.shapes.add_textbox(x + Inches(0.3), y + Inches(0.9), card_width - Inches(0.6), Inches(1.2))
        tf = features_box.text_frame
        tf.word_wrap = True

        for j, feature in enumerate(sec['features']):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = "• " + feature
            p.font.size = Pt(12)
            p.font.color.rgb = WHITE
            p.space_after = Pt(4)

    return slide

def add_comparison_slide(prs, title, left_title, left_points, right_title, right_points, section=None):
    return add_two_column_slide(prs, title, left_title, left_points, right_title, right_points, section)

def add_case_study_slide(prs, title, section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, DARK_NAVY)

    # Orange accent bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.2), prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = MUTED_ORANGE
    shape.line.fill.background()

    # Section label if provided
    if section:
        section_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(3), Inches(0.4))
        tf = section_box.text_frame
        p = tf.paragraphs[0]
        p.text = section.upper()
        p.font.size = Pt(12)
        p.font.color.rgb = MUTED_ORANGE
        p.font.bold = True
        p.letter_spacing = Pt(2)

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), prs.slide_width - Inches(1.6), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), prs.slide_width - Inches(1.6), Inches(0.4))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "When repeat business is systemized, growth becomes easier"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(180, 180, 190)

    # Case study cards
    card_width = Inches(5.5)
    card_height = Inches(3.5)
    gap = Inches(0.8)

    # Card 1
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1), Inches(2), card_width, card_height)
    card1.fill.solid()
    card1.fill.fore_color.rgb = RGBColor(30, 45, 67)
    card1.line.color.rgb = MUTED_ORANGE
    card1.line.width = Inches(0.03)

    # Card 1 content
    c1_title = slide.shapes.add_textbox(Inches(1.3), Inches(2.2), card_width - Inches(0.6), Inches(0.5))
    tf = c1_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Store A - Urban Location"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = MUTED_ORANGE

    c1_metrics = slide.shapes.add_textbox(Inches(1.3), Inches(2.8), card_width - Inches(0.6), Inches(2.5))
    tf = c1_metrics.text_frame
    metrics = [
        "↑ 45% repeat orders in 6 months",
        "↑ 60% pickup & delivery adoption",
        "↑ 30% faster payment collection",
        "↓ 50% manual follow-up time"
    ]
    for i, m in enumerate(metrics):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = m
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)

    # Card 2
    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2), card_width, card_height)
    card2.fill.solid()
    card2.fill.fore_color.rgb = RGBColor(30, 45, 67)
    card2.line.color.rgb = GOLD_ACCENT
    card2.line.width = Inches(0.03)

    # Card 2 content
    c2_title = slide.shapes.add_textbox(Inches(7.1), Inches(2.2), card_width - Inches(0.6), Inches(0.5))
    tf = c2_title.text_frame
    p = tf.paragraphs[0]
    p.text = "Multi-store B - 3 Locations"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD_ACCENT

    c2_metrics = slide.shapes.add_textbox(Inches(7.1), Inches(2.8), card_width - Inches(0.6), Inches(2.5))
    tf = c2_metrics.text_frame
    metrics2 = [
        "↑ 40% consolidated revenue growth",
        "↑ 85% cross-store visibility",
        "↑ 25% customer retention rate",
        "↓ 70% reconciliation time"
    ]
    for i, m in enumerate(metrics2):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = m
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.space_after = Pt(10)

    return slide


def main():
    print("Creating presentation...")
    prs = create_presentation()

    # Slide 1 - Title Slide
    add_title_slide(
        prs,
        "The Future of Dry Cleaning & Laundry Business in 2026",
        "How modern drycleaners can attract more customers, increase repeat business, and scale with control",
        "Jodhpur | 28 March 2026"
    )

    # Slide 2 - Opening Hook
    add_content_slide(
        prs,
        "Running a drycleaning business is getting harder",
        [
            "Customer expectations are rising",
            "Competition is increasing",
            "Manual operations break with growth",
            "Repeat business is harder than it looks",
            "Most stores still run without systems"
        ],
        section="THE CHALLENGE"
    )

    # Slide 3 - Market Shift
    add_content_slide(
        prs,
        "Customer behavior has changed",
        [
            "Customers discover businesses online",
            "They expect pickup and delivery convenience",
            "They want order updates",
            "They prefer digital payments",
            "They reward consistency and speed"
        ],
        section="THE SHIFT"
    )

    # Slide 4 - Reality Check
    add_content_slide(
        prs,
        "Walk-ins are no longer a growth strategy",
        [
            "Footfall is inconsistent",
            "Referrals are slow",
            "Competition is everywhere",
            "Google presence influences trust",
            "Visibility drives conversion"
        ],
        section="REALITY CHECK"
    )

    # Slide 5 - Big Opportunity
    add_framework_slide(
        prs,
        "The winners will not always be cheaper — they will be more convenient",
        [
            {"title": "Acquire", "desc": "Customers better"},
            {"title": "Serve", "desc": "Customers faster"},
            {"title": "Retain", "desc": "Customers longer"},
            {"title": "Collect", "desc": "Payments smarter"},
            {"title": "Scale", "desc": "With control"}
        ],
        section="OPPORTUNITY"
    )

    # Slide 6 - Customer Acquisition
    add_content_slide(
        prs,
        "More customers do not come by chance. They come by system.",
        [
            "Google visibility",
            "Ratings and reviews",
            "Meta ads / local promotion",
            "WhatsApp enquiry response",
            "Campaign-led acquisition",
            "Follow-up discipline"
        ],
        section="ACQUISITION"
    )

    # Slide 7 - Discovery & Trust
    add_content_slide(
        prs,
        "Your next customer may discover you before they ever visit you",
        [
            "Google search matters",
            "Reviews build trust",
            "Social proof improves conversion",
            "Fast response wins leads"
        ],
        section="DISCOVERY"
    )

    # Slide 8 - Pickup & Delivery
    add_content_slide(
        prs,
        "Convenience is now part of the service",
        [
            "Customers value time",
            "Pickup increases frequency",
            "Delivery improves retention",
            "Manual coordination creates chaos"
        ],
        section="CONVENIENCE"
    )

    # Slide 9 - Omni-channel Experience
    add_content_slide(
        prs,
        "Customers want to reach you from any channel — and get one seamless experience",
        [
            "Store",
            "WhatsApp",
            "Phone",
            "Website",
            "App",
            "Payment link / QR"
        ],
        section="OMNI-CHANNEL"
    )

    # Slide 10 - Service Visibility
    add_content_slide(
        prs,
        "Updates build trust",
        [
            "Pickup confirmed",
            "Garments received",
            "Processing status",
            "Ready for delivery",
            "Payment pending / completed"
        ],
        section="VISIBILITY"
    )

    # Slide 11 - Retention Problem
    add_content_slide(
        prs,
        "Most drycleaners do not have a customer problem. They have a repeat-business problem.",
        [
            "One-time orders",
            "Weak follow-up",
            "No return incentives",
            "No visibility on drop-off",
            "No habit-building system"
        ],
        section="RETENTION"
    )

    # Slide 12 - Loyalty & Predictable Revenue
    add_content_slide(
        prs,
        "Predictable revenue beats unpredictable footfall",
        [
            "Loyalty points",
            "Prepaid packages",
            "Subscription plans",
            "Repeat campaigns",
            "Win-back communication"
        ],
        section="REVENUE"
    )

    # Slide 13 - Proof / Business Impact
    add_case_study_slide(
        prs,
        "When repeat business is systemized, growth becomes easier",
        section="PROOF"
    )

    # Slide 14 - Digital Payments
    add_content_slide(
        prs,
        "Revenue is not real until it is collected",
        [
            "Pending dues hurt cash flow",
            "Manual follow-up is inconsistent",
            "Payment links reduce friction",
            "Faster collection improves business confidence"
        ],
        section="PAYMENTS"
    )

    # Slide 15 - Collection Discipline
    add_content_slide(
        prs,
        "Cash flow improves when payment follow-up becomes a process",
        [
            "Automated reminders",
            "Payment links",
            "Store-wise visibility",
            "Reconciliation support",
            "Fewer missed payments"
        ],
        section="COLLECTION"
    )

    # Slide 16 - Multi-store Challenge
    add_comparison_slide(
        prs,
        "Growth creates complexity faster than it creates control",
        "Manual Operations",
        [
            "Staff inconsistency",
            "Weak process control",
            "Fragmented customer history",
            "Unreliable reporting",
            "Poor visibility across outlets"
        ],
        "Digital Operations",
        [
            "Standardized workflows",
            "Centralized control",
            "Unified customer view",
            "Real-time reporting",
            "Full visibility"
        ],
        section="MULTI-STORE"
    )

    # Slide 17 - Centralized Control
    add_content_slide(
        prs,
        "Multi-store growth works only when control is centralized",
        [
            "Standard business rules",
            "Central reporting",
            "Access control",
            "Offer / pricing consistency",
            "Customer and order visibility",
            "Better accountability"
        ],
        section="CENTRALIZED CONTROL"
    )

    # Slide 18 - Franchise Opportunity
    add_content_slide(
        prs,
        "Strong operators can become strong brands",
        [
            "Expansion needs repeatability",
            "Brands need standardization",
            "Franchises need visibility and control",
            "Systems make scale safer"
        ],
        section="FRANCHISE"
    )

    # Slide 19 - Franchise Control Framework
    add_content_slide(
        prs,
        "Franchise growth fails when systems stay manual",
        [
            "HQ visibility",
            "Store-level control",
            "Standard service setup",
            "Brand consistency",
            "Royalty / payment discipline",
            "Central reporting"
        ],
        section="FRANCHISE CONTROL"
    )

    # Slide 20 - Modern Business Framework (Hexagon)
    add_hexagon_slide(
        prs,
        "The modern garment care business runs on 6 engines",
        [
            "Customer Acquisition",
            "Pickup & Delivery",
            "Communication",
            "Retention & Loyalty",
            "Payments & Collections",
            "Multi-store / Franchise Control"
        ],
        section="FRAMEWORK"
    )

    # Slide 21 - Why Systems Matter
    add_comparison_slide(
        prs,
        "Manual effort cannot scale modern service expectations",
        "Manual Chaos",
        [
            "Missed follow-ups",
            "Delayed updates",
            "Collection gaps",
            "Store inconsistency",
            "Weak reporting",
            "Poor scale readiness"
        ],
        "System-Led Control",
        [
            "Automated follow-ups",
            "Real-time updates",
            "Payment tracking",
            "Consistent operations",
            "Data-driven decisions",
            "Scale-ready architecture"
        ],
        section="SYSTEMS"
    )

    # Slide 22 - QDC Outcome Positioning
    add_outcome_slide(
        prs,
        "QDC helps drycleaners build a stronger business — not just run software",
        [
            {"title": "Better customer handling", "desc": "Unified view across all channels"},
            {"title": "Pickup & delivery coordination", "desc": "End-to-end tracking & routing"},
            {"title": "Repeat business & loyalty", "desc": "Automated campaigns & rewards"},
            {"title": "Faster payment collection", "desc": "Links, reminders & reconciliation"},
            {"title": "Multi-store visibility", "desc": "HQ-level oversight & control"},
            {"title": "Franchise-level control", "desc": "Standardization at scale"}
        ],
        section="QDC OUTCOMES"
    )

    # Slide 23 - QDC Capability Snapshot
    add_ecosystem_slide(
        prs,
        "One connected platform for growth, service, and control",
        [
            {"title": "Store Operations", "features": ["Order management", "Garment tracking", "Service catalog"]},
            {"title": "Plant / Processing", "features": ["Process workflows", "Quality checks", "Turnaround tracking"]},
            {"title": "Customer Communication", "features": ["WhatsApp integration", "SMS updates", "App notifications"]},
            {"title": "Payments", "features": ["Payment links", "UPI/QR", "Pending tracking"]},
            {"title": "Loyalty & Growth", "features": ["Points system", "Packages", "Campaigns"]},
            {"title": "Multi-store Control", "features": ["HQ dashboard", "Store analytics", "Access control"]}
        ],
        section="CAPABILITIES"
    )

    # Slide 24 - Vision Close
    add_content_slide(
        prs,
        "The next generation of drycleaning businesses will be built differently",
        [
            "More visible",
            "More convenient",
            "More systemized",
            "More customer-driven",
            "More scalable"
        ],
        section="VISION"
    )

    # Slide 25 - Closing CTA
    add_cta_slide(
        prs,
        "The right time to modernize is before growth exposes the gaps",
        [
            "Book a business consultation",
            "Get a live demo",
            "Discuss your growth plan",
            "Explore how QDC fits"
        ],
        section="NEXT STEPS"
    )

    # Save presentation
    output_path = "/Users/ashwanikataria/Documents/slides/laundrylaunchpad/Laundry_Business_Presentation_2026.pptx"
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    main()